"""Generate NUMMF SIH 2026 Presentation PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import math
import os

# ── Palette ──────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1e, 0x3a, 0x5f)
DARK_BG   = RGBColor(0x0f, 0x17, 0x2a)
PURPLE    = RGBColor(0xd9, 0x46, 0xef)
EMERALD   = RGBColor(0x10, 0xb9, 0x81)
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)
WHITE     = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY= RGBColor(0x94, 0xa3, 0xb8)
MID_GRAY  = RGBColor(0x47, 0x56, 0x69)
PANEL_BG  = RGBColor(0x1e, 0x29, 0x3b)
ACCENT2   = RGBColor(0x3b, 0x82, 0xf6)   # blue

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ═══════════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════════
def add_rect(slide, left, top, width, height, fill=None, line=None, line_width=Pt(0)):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri", italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.name      = font_name
    run.font.italic    = italic
    return txb


def add_paragraph(tf, text, font_size=14, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, space_before=Pt(4), font_name="Calibri", italic=False):
    p   = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.name      = font_name
    run.font.italic    = italic
    return p


def dark_slide(slide):
    """Fill slide with dark background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def header_bar(slide, title_text):
    """Top nav bar with title."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.65), fill=NAVY)
    add_textbox(slide, Inches(0.3), Inches(0.1), Inches(8), Inches(0.5),
                "NUMMF  |  " + title_text, font_size=12, color=LIGHT_GRAY, font_name="Calibri")
    add_textbox(slide, Inches(11.5), Inches(0.1), Inches(1.7), Inches(0.5),
                "SIH 2026 · ID 26099", font_size=10, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, tagline="One Nation, One Material Code"):
    add_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), fill=NAVY)
    add_textbox(slide, Inches(0.3), Inches(7.13), Inches(10), Inches(0.25),
                tagline, font_size=9, color=MID_GRAY)
    add_textbox(slide, Inches(11), Inches(7.13), Inches(2), Inches(0.25),
                "NUMMF v1.0", font_size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def card(slide, left, top, width, height, title=None, body_lines=None,
         title_size=13, body_size=11, border_color=PURPLE, fill=RGBColor(0x11,0x18,0x27)):
    """Rounded-ish card (rectangle with subtle border)."""
    s = add_rect(slide, left, top, width, height, fill=fill,
                 line=border_color, line_width=Pt(1.2))
    if title:
        # Title strip
        add_rect(slide, left, top, width, Inches(0.35), fill=NAVY)
        add_textbox(slide, left+Inches(0.12), top+Inches(0.04),
                    width-Inches(0.24), Inches(0.3), title,
                    font_size=title_size, bold=True, color=PURPLE)
        body_top = top + Inches(0.38)
        body_h   = height - Inches(0.38)
    else:
        body_top = top + Inches(0.1)
        body_h   = height - Inches(0.2)

    if body_lines:
        txb = slide.shapes.add_textbox(left+Inches(0.12), body_top,
                                        width-Inches(0.24), body_h)
        tf  = txb.text_frame
        tf.word_wrap = True
        first = True
        for line in body_lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(2)
            parts = line.split("::")
            label = parts[0].strip() if parts else ""
            rest  = parts[1].strip() if len(parts) > 1 else ""
            if label and rest:
                r1 = p.add_run(); r1.text = label + " "
                r1.font.size = Pt(body_size); r1.font.bold = True; r1.font.color.rgb = LIGHT_GRAY; r1.font.name = "Calibri"
                r2 = p.add_run(); r2.text = rest
                r2.font.size = Pt(body_size); r2.font.color.rgb = WHITE; r2.font.name = "Calibri"
            else:
                r = p.add_run(); r.text = line
                r.font.size = Pt(body_size); r.font.color.rgb = WHITE; r.font.name = "Calibri"
    return s


def bullet_textbox(slide, left, top, width, height, bullets, font_size=13, color=WHITE):
    """A textbox with bulleted list."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = "  ▸  " + b
        r.font.size = Pt(font_size); r.font.color.rgb = color; r.font.name = "Calibri"
    return txb


def accent_line(slide, left, top, width=Inches(4), color=PURPLE, height=Pt(3)):
    add_rect(slide, left, top, width, height, fill=color)


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
def slide01_title():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    # Gradient-like stripes
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=PURPLE)
    add_rect(sl, Inches(0.18), Inches(0), Inches(0.08), SLIDE_H, fill=ACCENT2)
    # Top-right badge
    add_rect(sl, Inches(10.8), Inches(0.2), Inches(2.3), Inches(0.45), fill=NAVY, line=PURPLE, line_width=Pt(1))
    add_textbox(sl, Inches(10.85), Inches(0.25), Inches(2.2), Inches(0.35),
                "SIH 2026  ·  Problem ID 26099", font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT, font_name="Calibri")
    # Logo box
    add_rect(sl, Inches(0.4), Inches(0.3), Inches(1.8), Inches(1.2), fill=NAVY, line=PURPLE, line_width=Pt(1.5))
    add_textbox(sl, Inches(0.45), Inches(0.4), Inches(1.7), Inches(0.55), "NUMMF", font_size=22, bold=True, color=PURPLE, font_name="Calibri")
    add_textbox(sl, Inches(0.45), Inches(0.9), Inches(1.7), Inches(0.45), "National Unified\nMaterial Master\nFramework", font_size=9, color=LIGHT_GRAY, font_name="Calibri")
    # Center title
    add_textbox(sl, Inches(2.5), Inches(2.2), Inches(8.3), Inches(1.2),
                "One Nation,", font_size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_textbox(sl, Inches(2.5), Inches(3.2), Inches(8.3), Inches(1.0),
                "One Material Code", font_size=58, bold=True, color=PURPLE, align=PP_ALIGN.CENTER, font_name="Calibri")
    # Subtitle
    add_textbox(sl, Inches(3), Inches(4.4), Inches(7.3), Inches(0.7),
                "AI-driven standardization across material records from CPSEs",
                font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri", italic=True)
    # Tagline
    add_textbox(sl, Inches(3.5), Inches(5.1), Inches(6.3), Inches(0.5),
                "Normalize · Match · Assign · Harmonize",
                font_size=14, color=ACCENT2, align=PP_ALIGN.CENTER, font_name="Calibri")
    # CTA
    add_rect(sl, Inches(5.5), Inches(5.9), Inches(2.3), Inches(0.7), fill=PURPLE)
    add_textbox(sl, Inches(5.5), Inches(5.95), Inches(2.3), Inches(0.6),
                "▶   START DEMO", font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    # Tech badges
    badges = ["FastAPI", "React + TS", "sentence-transformers", "Docker", "PostgreSQL"]
    for i, b in enumerate(badges):
        x = Inches(2.8) + i * Inches(1.95)
        add_rect(sl, x, Inches(6.9), Inches(1.8), Inches(0.3), fill=NAVY, line=MID_GRAY, line_width=Pt(0.5))
        add_textbox(sl, x+Inches(0.05), Inches(6.92), Inches(1.7), Inches(0.25),
                    b, font_size=8, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")
    # Footer
    add_textbox(sl, Inches(0.3), Inches(7.2), Inches(4), Inches(0.25),
                "Ministry of Petroleum & Natural Gas  ·  CPCL", font_size=9, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════
def slide02_problem():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "The Problem")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(12.5), Inches(0.8),
                "CPSEs Face Fragmented Material Catalogs",
                font_size=32, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.5), Inches(4), PURPLE)
    # Three columns
    cols = [
        ("IOCL", "Hex Bolt M20x100\nSS304 Grade 8.8", "Code: IOCL-HEX-001234", ACCENT2),
        ("NTPC", "Hex Bolt M20x100\nSS304 Grade 8.8", "Code: NT-BOLT-5678", AMBER),
        ("SAIL", "Hexagonal Bolt 20mm\nSS304", "Code: SA-FA-9012", EMERALD),
    ]
    for i, (org, desc, code, col) in enumerate(cols):
        x = Inches(0.5) + i * Inches(4.2)
        add_rect(sl, x, Inches(2.0), Inches(3.8), Inches(3.5), fill=NAVY, line=col, line_width=Pt(1.5))
        add_textbox(sl, x+Inches(0.15), Inches(2.1), Inches(3.5), Inches(0.4),
                    org, font_size=18, bold=True, color=col, font_name="Calibri")
        add_textbox(sl, x+Inches(0.15), Inches(2.55), Inches(3.5), Inches(1.2),
                    desc, font_size=13, color=WHITE, font_name="Calibri")
        add_rect(sl, x+Inches(0.15), Inches(3.8), Inches(3.5), Inches(0.45), fill=col)
        add_textbox(sl, x+Inches(0.15), Inches(3.85), Inches(3.5), Inches(0.35),
                    code, font_size=11, bold=True, color=DARK_BG, font_name="Calibri")
    # Impact
    add_textbox(sl, Inches(0.4), Inches(5.85), Inches(3), Inches(0.4),
                "❌  Impact:", font_size=14, bold=True, color=AMBER, font_name="Calibri")
    impacts = ["Duplicate master data across CPSEs", "No demand aggregation in procurement",
               "Higher inventory carrying costs", "Slower sourcing & tender preparation"]
    for i, imp in enumerate(impacts):
        add_rect(sl, Inches(0.4) + i*Inches(3.15), Inches(6.3), Inches(3), Inches(0.55),
                 fill=RGBColor(0x1a,0x23,0x2e), line=RGBColor(0x2d,0x3a,0x4f), line_width=Pt(0.5))
        add_textbox(sl, Inches(0.5) + i*Inches(3.15), Inches(6.38), Inches(2.8), Inches(0.45),
                    imp, font_size=10, color=LIGHT_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════
def slide03_solution():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "The Solution")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(8), Inches(0.8),
                "NUMMF — AI-Powered Unified Material Master",
                font_size=30, bold=True, color=WHITE, font_name="Calibri")
    add_textbox(sl, Inches(0.4), Inches(1.45), Inches(12), Inches(0.4),
                "Use AI/NLP to find equivalent materials across CPSEs and generate a single Common National Material Code (CNMC)",
                font_size=12, color=LIGHT_GRAY, italic=True, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.9), Inches(5), PURPLE)
    # Flow steps
    steps = [
        ("1", "INGEST", "Connect & import\nmaterial masters", ACCENT2),
        ("2", "NORMALIZE", "Domain-aware\ntext normalization", PURPLE),
        ("3", "AI MATCH", "4-stage scoring\npipeline", EMERALD),
        ("4", "CNMC", "Assign unified\nmaterial codes", AMBER),
    ]
    for i, (num, title, subtitle, col) in enumerate(steps):
        x = Inches(0.5) + i * Inches(3.15)
        add_rect(sl, x, Inches(2.3), Inches(2.9), Inches(1.5), fill=NAVY, line=col, line_width=Pt(1.5))
        add_rect(sl, x, Inches(2.3), Inches(0.55), Inches(1.5), fill=col)
        add_textbox(sl, x+Inches(0.08), Inches(0.78), Inches(0.4), Inches(0.5),
                    num, font_size=20, bold=True, color=DARK_BG, font_name="Calibri")
        add_textbox(sl, x+Inches(0.65), Inches(2.38), Inches(2.1), Inches(0.4),
                    title, font_size=14, bold=True, color=col, font_name="Calibri")
        add_textbox(sl, x+Inches(0.65), Inches(2.78), Inches(2.1), Inches(0.9),
                    subtitle, font_size=11, color=LIGHT_GRAY, font_name="Calibri")
        if i < 3:
            add_textbox(sl, x+Inches(2.95), Inches(2.95), Inches(0.25), Inches(0.4),
                        "→", font_size=20, color=MID_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")
    # Key Actions
    actions = [
        "Ingest material data from multiple CPSE ERP/SAP systems",
        "Normalize descriptions to a standard format",
        "Match materials using 4-stage AI pipeline",
        "Generate unique CNMC codes (CNMC-{SEGMENT}-{HASH})",
        "Map CPSE codes ↔ CNMC codes bidirectionally",
        "Review & Approve AI suggestions with full audit trail",
    ]
    card(sl, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.7),
         title="▸  Key Actions", body_lines=actions, title_size=14, body_size=12, border_color=PURPLE)
    add_textbox(sl, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.35),
                "Tagline:  \"One Nation, One Material Code\"", font_size=13, bold=True,
                color=PURPLE, align=PP_ALIGN.CENTER, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 4 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════
def slide04_techstack():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Technology Stack")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(8), Inches(0.8),
                "Built for Scale, Speed & Accuracy", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.5), Inches(4), PURPLE)
    layers = [
        ("Backend",   ["FastAPI — Async, type-safe, OpenAPI", "SQLAlchemy 2.0 — ORM + migrations",
                        "Pydantic V2 — Data validation", "JWT Auth — Secure API access"], ACCENT2),
        ("AI / ML",   ["sentence-transformers — State-of-the-art embeddings",
                        "RapidFuzz — Fast fuzzy string matching",
                        "Cross-Encoder — High-accuracy re-ranking",
                        "scikit-learn — TF-IDF fallback"], PURPLE),
        ("Frontend",  ["React 18 + TypeScript — Modern UI", "Tailwind CSS — Utility-first styling",
                        "Zustand — Lightweight state management", "Recharts — Beautiful data viz"], EMERALD),
        ("Data",      ["SQLite (dev) / PostgreSQL (prod)", "Alembic — DB migrations",
                        "Audit trail — Full governance log"], AMBER),
        ("Deploy",    ["Docker Compose — One-command deploy",
                        "Nginx — Reverse proxy + static cache",
                        "Production-ready → K8s scale-out"], LIGHT_GRAY),
    ]
    for i, (layer, items, col) in enumerate(layers):
        x = Inches(0.4) + (i % 3) * Inches(4.25)
        y = Inches(2.0) + (i // 3) * Inches(2.5)
        add_rect(sl, x, y, Inches(4.0), Inches(2.2), fill=NAVY, line=col, line_width=Pt(1.5))
        add_rect(sl, x, y, Inches(4.0), Inches(0.38), fill=col)
        add_textbox(sl, x+Inches(0.12), y+Inches(0.05), Inches(3.8), Inches(0.3),
                    layer, font_size=13, bold=True, color=DARK_BG, font_name="Calibri")
        for j, item in enumerate(items):
            add_textbox(sl, x+Inches(0.15), y+Inches(0.48)+j*Inches(0.38),
                        Inches(3.7), Inches(0.35), "▸  " + item,
                        font_size=10, color=WHITE, font_name="Calibri")
    add_textbox(sl, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
                "★  All libraries are open-source — no proprietary AI APIs required",
                font_size=11, bold=True, color=EMERALD, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 5 — AI PIPELINE
# ═══════════════════════════════════════════════════════════════════════
def slide05_aipipeline():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "AI Matching Pipeline")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "4-Stage AI Matching Pipeline",
                font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    stages = [
        ("1", "NORMALIZE", "Input: Raw ERP descriptions\nRemoves IS codes, standard codes\nNormalizes grades, dimensions, UOMs",
         ACCENT2, "30%"),
        ("2", "LEXICAL (30%)", "Token sort ratio + fuzzy match\nRapidFuzz-powered scoring\nHandles spelling variants & typos",
         ACCENT2, "30%"),
        ("3", "SEMANTIC (40%)", "Bi-encoder: all-MiniLM-L6-v2\nDense vector embeddings\nCaptures meaning beyond words",
         PURPLE, "40%"),
        ("4", "NUMERIC (30%)", "Family, grade, dimension\nUOM compatibility check\nPhysical attribute matching",
         AMBER, "30%"),
        ("5", "RE-RANKING", "Cross-Encoder: ms-marco-MiniLM\nPrecision scoring on top-K\nFinal match classification",
         EMERALD, "Final"),
    ]
    for i, (num, title, desc, col, score) in enumerate(stages):
        x = Inches(0.35) + i * Inches(2.55)
        add_rect(sl, x, Inches(2.1), Inches(2.4), Inches(2.4), fill=NAVY, line=col, line_width=Pt(1.5))
        add_rect(sl, x, Inches(2.1), Inches(0.5), Inches(2.4), fill=col)
        add_textbox(sl, x+Inches(0.06), Inches(0.7), Inches(0.4), Inches(0.4),
                    num, font_size=18, bold=True, color=DARK_BG, font_name="Calibri")
        add_textbox(sl, x+Inches(0.6), Inches(2.15), Inches(1.7), Inches(0.4),
                    title, font_size=11, bold=True, color=col, font_name="Calibri")
        add_rect(sl, x+Inches(0.6), Inches(2.55), Inches(1.7), Inches(0.25), fill=RGBColor(0x0a,0x0e,0x17))
        add_textbox(sl, x+Inches(0.6), Inches(2.55), Inches(1.7), Inches(0.25),
                    score, font_size=10, bold=True, color=col, align=PP_ALIGN.CENTER, font_name="Calibri")
        add_textbox(sl, x+Inches(0.15), Inches(2.85), Inches(2.1), Inches(1.5),
                    desc, font_size=9, color=LIGHT_GRAY, font_name="Calibri")
        if i < 4:
            add_textbox(sl, x+Inches(2.45), Inches(3.0), Inches(0.15), Inches(0.4),
                        "→", font_size=20, color=MID_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")
    # Benefits
    benefits = [
        ("Speed", "Bi-encoder for fast candidate selection (ms latency)"),
        ("Accuracy", "Cross-encoder for precise re-ranking (top-K only)"),
        ("Domain-aware", "Normalizer handles 100+ Indian industrial terms"),
    ]
    for i, (b_title, b_body) in enumerate(benefits):
        x = Inches(0.4) + i * Inches(4.3)
        add_rect(sl, x, Inches(4.8), Inches(4.0), Inches(1.5), fill=NAVY, line=PURPLE, line_width=Pt(1))
        add_textbox(sl, x+Inches(0.15), Inches(4.9), Inches(3.7), Inches(0.35),
                    "✦  " + b_title, font_size=13, bold=True, color=PURPLE, font_name="Calibri")
        add_textbox(sl, x+Inches(0.15), Inches(5.3), Inches(3.7), Inches(0.8),
                    b_body, font_size=11, color=LIGHT_GRAY, font_name="Calibri")
    # Threshold
    add_rect(sl, Inches(4.8), Inches(6.5), Inches(3.7), Inches(0.5), fill=EMERALD)
    add_textbox(sl, Inches(4.8), Inches(6.53), Inches(3.7), Inches(0.4),
                "Threshold: ≥65% composite score → submit for human review",
                font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 6 — DOMAIN INTELLIGENCE
# ═══════════════════════════════════════════════════════════════════════
def slide06_domain():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Domain Intelligence")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "Indian Industrial Material Knowledge",
                font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    # Normalizer table
    add_rect(sl, Inches(0.4), Inches(2.0), Inches(5.8), Inches(3.8), fill=NAVY, line=PURPLE, line_width=Pt(1))
    add_rect(sl, Inches(0.4), Inches(2.0), Inches(5.8), Inches(0.38), fill=PURPLE)
    add_textbox(sl, Inches(0.5), Inches(2.05), Inches(5.6), Inches(0.3),
                "Normalizer Features", font_size=13, bold=True, color=DARK_BG, font_name="Calibri")
    features = [
        ("IS Code removal",   '"IS:1239 Pipe 50NB"',          '"pipe 50NB"'),
        ("Grade normalization",   '"SS304, SS-304, 304SS"',      '"SS304"'),
        ("Dimension extraction",   '"M20x100"',                   "{thread: 20, length: 100}"),
        ("UOM suggestion",    '"Hydraulic Oil 200L"',          '"LTR"'),
        ("Variant mapping",   '"Flange, Flg., FLNG"',          '"flange"'),
        ("Material family",   '"Hex Bolt" → FA segment lookup', 'Auto-classified'),
    ]
    headers = ["Feature", "Example Input", "Normalized Output"]
    for ci, h in enumerate(headers):
        add_textbox(sl, Inches(0.5)+ci*Inches(1.9), Inches(2.45), Inches(1.8), Inches(0.3),
                    h, font_size=10, bold=True, color=LIGHT_GRAY, font_name="Calibri")
    for ri, (feat, inp, out) in enumerate(features):
        y = Inches(2.8) + ri * Inches(0.42)
        add_rect(sl, Inches(0.4), y, Inches(5.8), Inches(0.4),
                 fill=RGBColor(0x14,0x1c,0x29) if ri%2==0 else RGBColor(0x18,0x22,0x32),
                 line=RGBColor(0x2d,0x3a,0x4f), line_width=Pt(0.3))
        add_textbox(sl, Inches(0.5), y+Inches(0.05), Inches(1.85), Inches(0.3), feat, font_size=9, color=LIGHT_GRAY, font_name="Calibri")
        add_textbox(sl, Inches(2.4), y+Inches(0.05), Inches(1.85), Inches(0.3), inp, font_size=9, color=AMBER, font_name="Calibri")
        add_textbox(sl, Inches(4.3), y+Inches(0.05), Inches(1.85), Inches(0.3), out, font_size=9, color=EMERALD, font_name="Calibri")
    # Families
    add_rect(sl, Inches(6.5), Inches(2.0), Inches(6.4), Inches(3.8), fill=NAVY, line=ACCENT2, line_width=Pt(1))
    add_rect(sl, Inches(6.5), Inches(2.0), Inches(6.4), Inches(0.38), fill=ACCENT2)
    add_textbox(sl, Inches(6.6), Inches(2.05), Inches(6.2), Inches(0.3),
                "14 Material Families (Segments)", font_size=13, bold=True, color=DARK_BG, font_name="Calibri")
    families = [
        ("FA", "Fasteners"), ("PT", "Pipes & Tubes"), ("VF", "Valves & Fittings"),
        ("EL", "Electrical"), ("BE", "Bearings"), ("HL", "Hydraulics & Pneumatics"),
        ("IN", "Instruments"), ("SS", "Structural Steel"), ("WE", "Welding"),
        ("SP", "Safety & PPE"), ("PC", "Pumps & Compressors"), ("CH", "Chemicals"),
        ("PK", "Packaging"), ("MX", "Miscellaneous"),
    ]
    for i, (code, name) in enumerate(families):
        col = i % 2
        row = i // 2
        x = Inches(6.6) + col * Inches(3.1)
        y = Inches(2.5) + row * Inches(0.42)
        add_rect(sl, x, y, Inches(3.0), Inches(0.36),
                 fill=RGBColor(0x14,0x1c,0x29), line=PURPLE, line_width=Pt(0.4))
        add_rect(sl, x, y, Inches(0.55), Inches(0.36), fill=PURPLE)
        add_textbox(sl, x+Inches(0.06), y+Inches(0.04), Inches(0.45), Inches(0.28),
                    code, font_size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER, font_name="Calibri")
        add_textbox(sl, x+Inches(0.62), y+Inches(0.04), Inches(2.3), Inches(0.28),
                    name, font_size=10, color=WHITE, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 7 — CNMC CODES
# ═══════════════════════════════════════════════════════════════════════
def slide07_cnmc():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "CNMC Code Generation")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "Common National Material Code (CNMC)",
                font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    add_textbox(sl, Inches(0.4), Inches(1.55), Inches(12), Inches(0.35),
                "Format:  CNMC-{SEGMENT}-{HASH}     |     Stable · Unique · Human-readable · Traceable",
                font_size=13, color=LIGHT_GRAY, font_name="Calibri")
    # Example flow
    flow = [
        ("Input Material", '"Hex Bolt M20x100\nSS304 Grade 8.8"', ACCENT2),
        ("Segment Lookup", "Segment: FA\n(Fasteners)", PURPLE),
        ("Semantic Hash", 'MD5("hex bolt M20x100\nSS304")[:6] = "4A2B9C"', EMERALD),
        ("CNMC Code", "CNMC-FA-4A2B9C", AMBER),
    ]
    for i, (label, text, col) in enumerate(flow):
        x = Inches(0.4) + i * Inches(3.15)
        y = Inches(2.2)
        add_rect(sl, x, y, Inches(2.9), Inches(1.5), fill=NAVY, line=col, line_width=Pt(1.5))
        add_rect(sl, x, y, Inches(2.9), Inches(0.38), fill=col)
        add_textbox(sl, x+Inches(0.12), y+Inches(0.05), Inches(2.7), Inches(0.3),
                    label, font_size=11, bold=True, color=DARK_BG, font_name="Calibri")
        add_textbox(sl, x+Inches(0.15), y+Inches(0.48), Inches(2.6), Inches(0.9),
                    text, font_size=11, color=WHITE, font_name="Calibri")
        if i < 3:
            add_textbox(sl, x+Inches(3.0), y+Inches(0.45), Inches(0.2), Inches(0.4),
                        "→", font_size=22, color=MID_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")
    # Benefits
    benefits = [
        ("Stable", "Same material → same code, always reproducible"),
        ("Unique", "MD5 hash prevents collisions across millions of materials"),
        ("Readable", "FA prefix instantly tells you it's a fastener"),
        ("Traceable", "Links back to original CPSE codes in audit log"),
    ]
    for i, (b, desc) in enumerate(benefits):
        x = Inches(0.4) + i * Inches(3.2)
        add_rect(sl, x, Inches(4.0), Inches(3.0), Inches(1.5), fill=NAVY, line=EMERALD, line_width=Pt(1))
        add_textbox(sl, x+Inches(0.15), Inches(4.1), Inches(2.7), Inches(0.35),
                    "✦  " + b, font_size=13, bold=True, color=EMERALD, font_name="Calibri")
        add_textbox(sl, x+Inches(0.15), Inches(4.5), Inches(2.7), Inches(0.8),
                    desc, font_size=10, color=LIGHT_GRAY, font_name="Calibri")
    # Mapping table
    add_rect(sl, Inches(0.4), Inches(5.75), Inches(12.5), Inches(1.2), fill=NAVY, line=PURPLE, line_width=Pt(1))
    add_textbox(sl, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.3),
                "Bidirectional Mapping:  IOCL-HEX-001234  ↔  NT-BOLT-5678  ↔  SA-FA-9012  →  CNMC-FA-4A2B9C",
                font_size=11, bold=True, color=PURPLE, font_name="Calibri")
    add_textbox(sl, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.6),
                "Every CPSE code is permanently linked to its CNMC. Change source data? Mapping auto-updates. Full history in audit log.",
                font_size=10, color=LIGHT_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 8 — DASHBOARD & FEATURES
# ═══════════════════════════════════════════════════════════════════════
def slide08_dashboard():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Dashboard & Features")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "Complete Platform Features", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    features = [
        ("📊  Dashboard",    "Real-time KPIs, charts, CPSE breakdown\nDuplicate distribution, quality metrics"),
        ("🔍  AI Matching",  "4-stage pipeline execution\nCross-CPSE & intra-CPSE matching"),
        ("📦  Materials",    "Search, filter, browse all materials\nFamily & sub-family grouping"),
        ("✓  Review Queue",  "Approve/reject AI suggestions\nBulk actions, confidence filters"),
        ("🏷️  CNMC Codes",   "Auto-assign unified codes\nBatch generation, validation"),
        ("📈  Analytics",    "Quality scores, procurement savings\nTrend analysis, export reports"),
        ("⚙️  Admin",        "CPSE management, demo data seeding\nUser management, audit trail"),
        ("🔒  Audit Log",    "Full governance trail\nWho changed what, when, why"),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 4
        row = i // 4
        x = Inches(0.35) + col * Inches(3.2)
        y = Inches(2.0) + row * Inches(2.5)
        add_rect(sl, x, y, Inches(3.05), Inches(2.25), fill=NAVY,
                 line=[ACCENT2, PURPLE, EMERALD, AMBER][col], line_width=Pt(1))
        add_textbox(sl, x+Inches(0.12), y+Inches(0.1), Inches(2.8), Inches(0.45),
                    title, font_size=13, bold=True, color=WHITE, font_name="Calibri")
        add_textbox(sl, x+Inches(0.12), y+Inches(0.6), Inches(2.8), Inches(1.4),
                    desc, font_size=10, color=LIGHT_GRAY, font_name="Calibri")
    # KPI strip
    kpis = [("326K+", "Total Rows"), ("51K+", "Duplicates"), ("90+", "Cross Matches"), ("48", "CNMC Codes")]
    for i, (val, label) in enumerate(kpis):
        x = Inches(0.35) + i * Inches(3.2)
        add_rect(sl, x, Inches(6.75), Inches(3.05), Inches(0.5), fill=EMERALD)
        add_textbox(sl, x+Inches(0.1), Inches(6.78), Inches(1.4), Inches(0.4),
                    val, font_size=14, bold=True, color=DARK_BG, font_name="Calibri")
        add_textbox(sl, x+Inches(1.5), Inches(6.8), Inches(1.5), Inches(0.35),
                    label, font_size=9, color=DARK_BG, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 9 — USER FLOW
# ═══════════════════════════════════════════════════════════════════════
def slide09_flow():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "How It Works")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "End-to-End User Journey", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    steps = [
        ("1", "CPSE Admin uploads\nmaterial data", "CSV / Excel / SAP IDoc", ACCENT2),
        ("2", "System normalizes\ndescriptions", "Domain normalizer active", PURPLE),
        ("3", "AI runs matching\npipeline across CPSEs", "4-stage: lexical + semantic + numeric", EMERALD),
        ("4", "Proposals generated\n(score ≥ 65%)", "Match proposals with confidence", AMBER),
        ("5", "User reviews\nproposals in queue", "Approve / Reject / Modify", ACCENT2),
        ("6", "CNMC assigned\nautomatically", "Stable code generated", PURPLE),
        ("7", "All CPSEs mapped\nto same CNMC", "Unified master record", EMERALD),
        ("8", "Analytics updated\nin real-time", "Dashboards, reports, KPIs", AMBER),
    ]
    for i, (num, title, sub, col) in enumerate(steps):
        row = i // 4
        col_i = i % 4
        x = Inches(0.4) + col_i * Inches(3.15)
        y = Inches(2.0) + row * Inches(2.4)
        add_rect(sl, x, y, Inches(2.95), Inches(2.1), fill=NAVY, line=col, line_width=Pt(1.5))
        add_rect(sl, x, y, Inches(0.5), Inches(2.1), fill=col)
        add_textbox(sl, x+Inches(0.07), y+Inches(0.02), Inches(0.38), Inches(0.4),
                    num, font_size=18, bold=True, color=DARK_BG, font_name="Calibri")
        add_textbox(sl, x+Inches(0.6), y+Inches(0.1), Inches(2.2), Inches(0.6),
                    title, font_size=11, bold=True, color=WHITE, font_name="Calibri")
        add_textbox(sl, x+Inches(0.6), y+Inches(0.75), Inches(2.2), Inches(0.5),
                    sub, font_size=9, color=LIGHT_GRAY, font_name="Calibri")
        # Arrow
        if i < 7:
            next_i = i + 1
            if next_i % 4 != 0:
                ax = x + Inches(3.0)
                ay = y + Inches(0.9)
                add_textbox(sl, ax, ay, Inches(0.2), Inches(0.4), "→", font_size=18, color=MID_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 10 — DIFFERENTIATORS
# ═══════════════════════════════════════════════════════════════════════
def slide10_diff():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Innovation & Differentiation")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "What Makes NUMMF Unique", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    # Table
    headers = ["Aspect", "Typical Solutions", "NUMMF"]
    rows = [
        ("Matching Method",     "Fuzzy string only",                 "4-stage AI pipeline"),
        ("Normalization",       "Basic lowercase / trim",             "100+ Indian industrial terms"),
        ("Code Generation",     "Random / UUID",                      "Semantic hash + segment prefix"),
        ("Cross-CPSE Matching", "No inter-org matching",              "Full multi-org matching"),
        ("Re-ranking",          "None",                               "Cross-encoder precision scoring"),
        ("Tech Dependencies",   "Proprietary AI APIs",               "100% open-source libraries"),
        ("Language Support",    "English only",                       "Hindi/Tamil/Telugu ready"),
        ("Governance",          "Basic logging",                      "Full audit trail + versioning"),
    ]
    col_widths = [Inches(3.0), Inches(4.0), Inches(4.8)]
    col_colors = [PURPLE, MID_GRAY, EMERALD]
    # Header
    for ci, h in enumerate(headers):
        x = Inches(0.4) + sum(col_widths[:ci])
        add_rect(sl, x, Inches(2.0), col_widths[ci], Inches(0.45), fill=NAVY, line=WHITE, line_width=Pt(1))
        add_textbox(sl, x+Inches(0.1), Inches(2.05), col_widths[ci]-Inches(0.2), Inches(0.35),
                    h, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    # Rows
    for ri, row in enumerate(rows):
        y = Inches(2.5) + ri * Inches(0.5)
        fill_bg = RGBColor(0x14,0x1c,0x29) if ri%2==0 else RGBColor(0x18,0x22,0x32)
        for ci, cell in enumerate(row):
            x = Inches(0.4) + sum(col_widths[:ci])
            add_rect(sl, x, y, col_widths[ci], Inches(0.48), fill=fill_bg,
                     line=RGBColor(0x2d,0x3a,0x4f), line_width=Pt(0.5))
            c = [PURPLE, MID_GRAY, EMERALD][ci]
            bold = (ci == 2)
            add_textbox(sl, x+Inches(0.1), y+Inches(0.08), col_widths[ci]-Inches(0.2), Inches(0.35),
                        cell, font_size=10, bold=bold, color=c, font_name="Calibri")
    # Novel contributions
    add_textbox(sl, Inches(0.4), Inches(6.7), Inches(5), Inches(0.3),
                "✦  Novel Contributions:", font_size=12, bold=True, color=PURPLE, font_name="Calibri")
    contribs = [
        "First Indian-industrial normalizer for material master data",
        "Hybrid bi-encoder + cross-encoder pipeline for speed + accuracy",
        "Stable semantic hashing for CNMC code generation",
        "Complete governance + audit trail for regulatory compliance",
    ]
    for i, c in enumerate(contribs):
        add_textbox(sl, Inches(0.4)+i*Inches(3.15), Inches(7.0), Inches(3.0), Inches(0.35),
                    "▸  " + c, font_size=9, color=LIGHT_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 11 — TECHNICAL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════
def slide11_arch():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Technical Architecture")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "System Architecture", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    # Architecture boxes
    comps = [
        ("Frontend", "React 18 + TS\nTailwind CSS\nZustand + Recharts", ACCENT2, Inches(0.4), Inches(2.0), Inches(2.8), Inches(2.0)),
        ("Backend API", "FastAPI (async)\n20+ REST endpoints\nJWT Auth", PURPLE, Inches(3.5), Inches(2.0), Inches(2.8), Inches(2.0)),
        ("AI Engine", "sentence-transformers\nRapidFuzz\nCross-Encoder", EMERALD, Inches(6.6), Inches(2.0), Inches(2.8), Inches(2.0)),
        ("Database", "SQLAlchemy 2.0\nSQLite / PostgreSQL\nAlembic migrations", AMBER, Inches(9.7), Inches(2.0), Inches(2.9), Inches(2.0)),
        ("Docker", "docker-compose up\nNginx reverse proxy\nFull stack deploy", LIGHT_GRAY, Inches(5.0), Inches(4.5), Inches(3.5), Inches(1.4)),
    ]
    for (title, desc, col, x, y, w, h) in comps:
        add_rect(sl, x, y, w, h, fill=NAVY, line=col, line_width=Pt(1.5))
        add_rect(sl, x, y, w, Inches(0.38), fill=col)
        add_textbox(sl, x+Inches(0.12), y+Inches(0.05), w-Inches(0.24), Inches(0.3),
                    title, font_size=12, bold=True, color=DARK_BG, font_name="Calibri")
        add_textbox(sl, x+Inches(0.12), y+Inches(0.48), w-Inches(0.24), h-Inches(0.55),
                    desc, font_size=10, color=WHITE, font_name="Calibri")
    # Arrows
    arrows = [
        (Inches(3.2), Inches(2.9), Inches(3.5), Inches(2.9)),
        (Inches(6.3), Inches(2.9), Inches(6.6), Inches(2.9)),
        (Inches(9.4), Inches(2.9), Inches(9.7), Inches(2.9)),
    ]
    for ax, ay, bx, by in arrows:
        add_textbox(sl, ax, ay, bx-ax, Inches(0.4), "→", font_size=20, color=MID_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")
    # API endpoints
    eps = ["Auth", "Materials", "Matching", "CNMC Mapping", "Analytics", "Admin"]
    for i, ep in enumerate(eps):
        x = Inches(0.4) + i * Inches(2.15)
        add_rect(sl, x, Inches(6.5), Inches(2.0), Inches(0.4),
                 fill=RGBColor(0x14,0x1c,0x29), line=MID_GRAY, line_width=Pt(0.5))
        add_textbox(sl, x+Inches(0.08), Inches(6.53), Inches(1.85), Inches(0.3),
                    ep, font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 12 — SCALABILITY
# ═══════════════════════════════════════════════════════════════════════
def slide12_scale():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Scalability & Deployment")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "Built to Scale", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    stages = [
        ("Development", "Single Docker", "SQLite", "CPU", ACCENT2),
        ("Staging", "Docker Compose", "PostgreSQL", "GPU (optional)", PURPLE),
        ("Production", "Docker + Nginx + K8s", "PostgreSQL RDS", "GPU cluster", EMERALD),
    ]
    headers = ["Stage", "Architecture", "Database", "ML Models"]
    col_w = [Inches(2.5), Inches(3.5), Inches(3.0), Inches(3.0)]
    for ci, h in enumerate(headers):
        x = Inches(0.4) + sum(col_w[:ci])
        add_rect(sl, x, Inches(2.0), col_w[ci], Inches(0.4), fill=NAVY, line=WHITE, line_width=Pt(1))
        add_textbox(sl, x+Inches(0.1), Inches(2.03), col_w[ci]-Inches(0.2), Inches(0.35),
                    h, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    for ri, (stage, arch, db, ml, col) in enumerate(stages):
        y = Inches(2.45) + ri * Inches(0.75)
        add_rect(sl, Inches(0.4), y, Inches(12.0), Inches(0.7), fill=RGBColor(0x14,0x1c,0x29), line=col, line_width=Pt(0.8))
        add_textbox(sl, Inches(0.5), y+Inches(0.12), Inches(2.3), Inches(0.45), stage, font_size=11, bold=True, color=col, font_name="Calibri")
        add_textbox(sl, Inches(2.95), y+Inches(0.12), Inches(3.3), Inches(0.45), arch, font_size=11, color=WHITE, font_name="Calibri")
        add_textbox(sl, Inches(6.3), y+Inches(0.12), Inches(2.8), Inches(0.45), db, font_size=11, color=WHITE, font_name="Calibri")
        add_textbox(sl, Inches(9.15), y+Inches(0.12), Inches(3.1), Inches(0.45), ml, font_size=11, color=WHITE, font_name="Calibri")
    # Scaling points
    scale_points = [
        ("Frontend", "CDN + Nginx cache → sub-second TTFB"),
        ("Backend", "Uvicorn workers (4+), load balanced"),
        ("AI Engine", "Bi-encoder pre-computes embeddings once, O(1) lookup"),
        ("Database", "Read replicas for analytics queries"),
    ]
    for i, (comp, desc) in enumerate(scale_points):
        x = Inches(0.4) + i * Inches(3.2)
        add_rect(sl, x, Inches(4.55), Inches(3.05), Inches(1.5), fill=NAVY, line=EMERALD, line_width=Pt(1))
        add_textbox(sl, x+Inches(0.12), Inches(4.65), Inches(2.8), Inches(0.35),
                    "⚡  " + comp, font_size=12, bold=True, color=EMERALD, font_name="Calibri")
        add_textbox(sl, x+Inches(0.12), Inches(5.05), Inches(2.8), Inches(0.8),
                    desc, font_size=10, color=LIGHT_GRAY, font_name="Calibri")
    # Command
    add_rect(sl, Inches(3.5), Inches(6.3), Inches(6.3), Inches(0.55), fill=RGBColor(0x0a,0x0e,0x17), line=EMERALD, line_width=Pt(1.5))
    add_textbox(sl, Inches(3.55), Inches(6.35), Inches(6.2), Inches(0.45),
                "docker-compose up --build    →    Full stack deployed in 1 command",
                font_size=13, bold=True, color=EMERALD, align=PP_ALIGN.CENTER, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 13 — FUTURE ROADMAP
# ═══════════════════════════════════════════════════════════════════════
def slide13_roadmap():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Future Roadmap")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "Beyond the Prototype", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    phases = [
        ("Phase 1\n(Now)", "Core platform + AI matching\nDocker Compose deploy\nDemo data & UI", EMERALD),
        ("Phase 2", "SAP integration (IDoc adapter)\nReal-time ERP sync\nMulti-CPSE live connection", ACCENT2),
        ("Phase 3", "Multi-language NLP\nHindi / Tamil / Telugu\nRegional material descriptions", PURPLE),
        ("Phase 4", "Knowledge graph\nMaterial taxonomy\nOntology-based matching", AMBER),
        ("Phase 5", "ML procurement savings\nAI-driven sourcing\nDemand forecasting", LIGHT_GRAY),
        ("Phase 6+", "Mobile app for field verification\nBlockchain audit trail\nStrategic sourcing AI", MID_GRAY),
    ]
    for i, (phase, desc, col) in enumerate(phases):
        x = Inches(0.4) + i * Inches(2.15)
        add_rect(sl, x, Inches(2.0), Inches(2.0), Inches(1.2), fill=NAVY, line=col, line_width=Pt(1.5))
        add_rect(sl, x, Inches(2.0), Inches(2.0), Inches(0.45), fill=col)
        add_textbox(sl, x+Inches(0.1), Inches(2.05), Inches(1.85), Inches(0.35),
                    phase, font_size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER, font_name="Calibri")
        add_textbox(sl, x+Inches(0.1), Inches(2.5), Inches(1.85), Inches(0.6),
                    desc, font_size=9, color=LIGHT_GRAY, font_name="Calibri")
        if i < 5:
            add_textbox(sl, x+Inches(2.05), Inches(2.55), Inches(0.15), Inches(0.4),
                        "→", font_size=18, color=MID_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")
    # Timeline bar
    add_rect(sl, Inches(0.4), Inches(3.5), Inches(12.5), Inches(0.08), fill=PURPLE)
    for i in range(7):
        x = Inches(0.4) + i * Inches(2.08)
        add_rect(sl, x, Inches(3.45), Inches(0.08), Inches(0.18), fill=PURPLE)
    # Details
    details = [
        ("Q3 2026", "MVP complete, SIH demo ready", EMERALD),
        ("Q4 2026", "Pilot with 2 CPSEs, SAP integration", ACCENT2),
        ("Q1 2027", "Expand to 5 CPSEs, multi-language", PURPLE),
        ("Q2 2027", "Knowledge graph, ML procurement AI", AMBER),
    ]
    for i, (q, desc, col) in enumerate(details):
        x = Inches(0.4) + i * Inches(3.2)
        y = Inches(3.8) + (i % 2) * Inches(0.9)
        add_rect(sl, x, y, Inches(3.0), Inches(0.75), fill=NAVY, line=col, line_width=Pt(1))
        add_textbox(sl, x+Inches(0.12), y+Inches(0.05), Inches(2.8), Inches(0.3),
                    q, font_size=11, bold=True, color=col, font_name="Calibri")
        add_textbox(sl, x+Inches(0.12), y+Inches(0.38), Inches(2.8), Inches(0.3),
                    desc, font_size=9, color=LIGHT_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 14 — IMPACT
# ═══════════════════════════════════════════════════════════════════════
def slide14_impact():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Impact & Benefits")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "Expected Impact", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    # CPSE benefits
    add_rect(sl, Inches(0.4), Inches(2.0), Inches(6.0), Inches(3.6), fill=NAVY, line=ACCENT2, line_width=Pt(1))
    add_rect(sl, Inches(0.4), Inches(2.0), Inches(6.0), Inches(0.38), fill=ACCENT2)
    add_textbox(sl, Inches(0.5), Inches(2.05), Inches(5.8), Inches(0.3),
                "For CPSEs", font_size=14, bold=True, color=DARK_BG, font_name="Calibri")
    cse_benefits = [
        "Unified material code across all CPSEs",
        "Reduction in duplicate master data",
        "Improved data quality scores",
        "Better inventory optimization",
        "Reduced procurement costs via demand aggregation",
    ]
    for i, b in enumerate(cse_benefits):
        add_textbox(sl, Inches(0.55), Inches(2.5)+i*Inches(0.5), Inches(5.7), Inches(0.4),
                    "▸  " + b, font_size=12, color=WHITE, font_name="Calibri")
    # National benefits
    add_rect(sl, Inches(6.8), Inches(2.0), Inches(6.1), Inches(3.6), fill=NAVY, line=PURPLE, line_width=Pt(1))
    add_rect(sl, Inches(6.8), Inches(2.0), Inches(6.1), Inches(0.38), fill=PURPLE)
    add_textbox(sl, Inches(6.9), Inches(2.05), Inches(5.9), Inches(0.3),
                "For the Nation", font_size=14, bold=True, color=DARK_BG, font_name="Calibri")
    nat_benefits = [
        "\"One Nation, One Material Code\" vision",
        "Common procurement platform for government",
        "Better inter-CPSE collaboration",
        "Foundation for strategic sourcing",
        "Transparency in public procurement",
    ]
    for i, b in enumerate(nat_benefits):
        add_textbox(sl, Inches(6.95), Inches(2.5)+i*Inches(0.5), Inches(5.8), Inches(0.4),
                    "▸  " + b, font_size=12, color=WHITE, font_name="Calibri")
    # ROI
    add_rect(sl, Inches(0.4), Inches(5.85), Inches(12.5), Inches(1.1), fill=EMERALD)
    add_textbox(sl, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.4),
                "ROI Estimate (per CPSE, 100K materials, 5% duplicate reduction)",
                font_size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_textbox(sl, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                "5,000 fewer entries  ×  ₹500/code maintenance  =  ₹25 Lakhs savings per CPSE  |  Scale to 100+ CPSEs = ₹2,500 Cr national savings",
                font_size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 15 — DEMO PLAN
# ═══════════════════════════════════════════════════════════════════════
def slide15_demo():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Demo Plan")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "Live Demo Plan (5 Minutes)", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    demos = [
        ("Demo 1: Dashboard",  "Show KPIs, pie charts, CPSE breakdown\nQuality metrics & trend lines", "0:00 – 0:45", ACCENT2),
        ("Demo 2: Matching",   "Run pipeline between IOCL ↔ NTPC\nWatch 4-stage progress in real-time", "0:45 – 1:30", PURPLE),
        ("Demo 3: Results",    "Show AI proposals with scores\nConfidence levels, match types, CNMC codes", "1:30 – 2:15", EMERALD),
        ("Demo 4: Review",     "Approve/reject proposals\nBulk actions, confidence filters", "2:15 – 3:00", AMBER),
        ("Demo 5: CNMC",       "Show CNMC codes auto-assigned\nBidirectional mapping table", "3:00 – 3:45", ACCENT2),
        ("Demo 6: Analytics",  "Deep dive into quality metrics\nProcurement savings calculator", "3:45 – 4:15", PURPLE),
        ("Demo 7: Admin",      "Seed demo data, view audit trail\nCPSE management", "4:15 – 4:45", EMERALD),
    ]
    for i, (title, desc, time, col) in enumerate(demos):
        y = Inches(2.0) + i * Inches(0.62)
        add_rect(sl, Inches(0.4), y, Inches(0.7), Inches(0.55), fill=col)
        add_textbox(sl, Inches(0.42), y+Inches(0.08), Inches(0.65), Inches(0.38),
                    str(i+1), font_size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER, font_name="Calibri")
        add_rect(sl, Inches(1.15), y, Inches(6.5), Inches(0.55), fill=NAVY, line=col, line_width=Pt(0.8))
        add_textbox(sl, Inches(1.25), y+Inches(0.04), Inches(6.3), Inches(0.25),
                    title, font_size=11, bold=True, color=WHITE, font_name="Calibri")
        add_textbox(sl, Inches(1.25), y+Inches(0.28), Inches(6.3), Inches(0.22),
                    desc, font_size=9, color=LIGHT_GRAY, font_name="Calibri")
        add_rect(sl, Inches(7.7), y, Inches(1.8), Inches(0.55), fill=RGBColor(0x14,0x1c,0x29), line=MID_GRAY, line_width=Pt(0.5))
        add_textbox(sl, Inches(7.75), y+Inches(0.12), Inches(1.7), Inches(0.3),
                    time, font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_textbox(sl, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.3),
                "Backup: Pre-recorded demo video + screenshots for any tech failures",
                font_size=10, color=AMBER, align=PP_ALIGN.CENTER, italic=True, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 16 — THANK YOU / TEAM
# ═══════════════════════════════════════════════════════════════════════
def slide16_team():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    dark_slide(sl)
    # Gradient-like accent
    add_rect(sl, Inches(0), Inches(0), Inches(0.15), SLIDE_H, fill=PURPLE)
    add_rect(sl, Inches(0.15), Inches(0), Inches(0.08), SLIDE_H, fill=ACCENT2)
    add_textbox(sl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2),
                "Thank You", font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_textbox(sl, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.8),
                "Built for Smart India Hackathon 2026", font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True, font_name="Calibri")
    add_textbox(sl, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                "\"One Nation, One Material Code\"", font_size=18, bold=True, color=PURPLE, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_rect(sl, Inches(4.5), Inches(4.1), Inches(4.3), Inches(0.05), fill=PURPLE)
    # Team
    roles = [
        ("AI/ML Lead", "sentence-transformers, matching pipeline"),
        ("Backend Lead", "FastAPI, SQLAlchemy, Docker"),
        ("Frontend Lead", "React, Tailwind, Recharts"),
        ("Database/DevOps", "PostgreSQL, migrations, deployment"),
        ("Product/Business", "Requirements, demo, documentation"),
    ]
    for i, (role, skill) in enumerate(roles):
        x = Inches(0.5) + (i % 3) * Inches(4.2)
        y = Inches(4.3) + (i // 3) * Inches(1.1)
        add_rect(sl, x, y, Inches(3.9), Inches(0.95), fill=NAVY, line=PURPLE, line_width=Pt(1))
        add_textbox(sl, x+Inches(0.15), y+Inches(0.08), Inches(3.6), Inches(0.3),
                    "[Team Member " + str(i+1) + "]", font_size=13, bold=True, color=WHITE, font_name="Calibri")
        add_textbox(sl, x+Inches(0.15), y+Inches(0.4), Inches(3.6), Inches(0.2),
                    role, font_size=10, bold=True, color=PURPLE, font_name="Calibri")
        add_textbox(sl, x+Inches(0.15), y+Inches(0.6), Inches(3.6), Inches(0.25),
                    skill, font_size=9, color=LIGHT_GRAY, font_name="Calibri")
    # Contact
    add_rect(sl, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.45), fill=NAVY)
    add_textbox(sl, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.35),
                "Email: [team@email.com]  |  GitHub: [github.com/team/nummf]  |  Demo: [nummf-demo.vercel.app]",
                font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE 17 — PROBLEM STATEMENT (extra context slide)
# ═══════════════════════════════════════════════════════════════════════
def slide17_problemstmt():
    sl = prs.slides.add_slide(BLANK)
    dark_slide(sl)
    header_bar(sl, "Problem Statement")
    footer(sl)
    add_textbox(sl, Inches(0.4), Inches(0.8), Inches(10), Inches(0.7),
                "SIH 2026 · Problem ID 26099", font_size=30, bold=True, color=WHITE, font_name="Calibri")
    accent_line(sl, Inches(0.4), Inches(1.45), Inches(5), PURPLE)
    add_rect(sl, Inches(0.4), Inches(2.0), Inches(12.5), Inches(3.8), fill=NAVY, line=PURPLE, line_width=Pt(1.5))
    add_rect(sl, Inches(0.4), Inches(2.0), Inches(12.5), Inches(0.4), fill=PURPLE)
    add_textbox(sl, Inches(0.5), Inches(2.05), Inches(12.3), Inches(0.3),
                "Problem Statement", font_size=13, bold=True, color=DARK_BG, font_name="Calibri")
    stmt = (
        "AI-Driven Standardization and Harmonization of Material Codes Across CPSEs\n\n"
        "Develop an AI-driven system to standardize and harmonize material codes across Central Public Sector "
        "Enterprises (CPSEs). The system should automatically map disparate material descriptions, grades, "
        "and specifications from different CPSE ERP systems to a unified Common National Material Code (CNMC). "
        "The solution must handle Indian industrial nomenclature, support batch processing of legacy data, "
        "and provide a human-in-the-loop review interface for expert validation."
    )
    add_textbox(sl, Inches(0.6), Inches(2.55), Inches(12.1), Inches(2.5),
                stmt, font_size=14, color=WHITE, font_name="Calibri")
    # Organization info
    add_rect(sl, Inches(0.4), Inches(6.0), Inches(12.5), Inches(1.0), fill=RGBColor(0x14,0x1c,0x29), line=MID_GRAY, line_width=Pt(0.5))
    add_textbox(sl, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.3),
                "Organization: Ministry of Petroleum & Natural Gas  |  Chennai Petroleum Corporation Limited (CPCL)",
                font_size=12, bold=True, color=ACCENT2, font_name="Calibri")
    add_textbox(sl, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.3),
                "Category: Software / Smart Automation  |  Team Size: Up to 6 members  |  Duration: 36 hours",
                font_size=11, color=LIGHT_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════
#  BUILD ALL SLIDES
# ═══════════════════════════════════════════════════════════════════════
slide01_title()
slide02_problem()
slide03_solution()
slide04_techstack()
slide05_aipipeline()
slide06_domain()
slide07_cnmc()
slide08_dashboard()
slide09_flow()
slide10_diff()
slide11_arch()
slide12_scale()
slide13_roadmap()
slide14_impact()
slide15_demo()
slide16_team()
slide17_problemstmt()

out = r"C:\Users\udayps\Desktop\sih\sih-2026-material-harmonization\docs\NUMMF_SIH2026_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
size = os.path.getsize(out)
print(f"Size: {size/1024:.1f} KB")
